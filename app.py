
import os, base64
import streamlit as st
import pandas as pd
from datasets import load_dataset
from sentence_transformers import SentenceTransformer, util
import docx2txt
from pptx import Presentation
import fitz 
from transformers import AutoTokenizer, AutoModelForSeq2SeqLM, pipeline
from langchain_community.vectorstores import FAISS
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_community.llms.huggingface_pipeline import HuggingFacePipeline
from langchain_core.prompts import PromptTemplate
from langchain_core.runnables import RunnablePassthrough
from langchain_core.output_parsers import StrOutputParser
import streamlit.components.v1 as components


RESUME_MATCHER_IMAGE_URL = "https://recruitryte.com/wp-content/uploads/2023/04/AI-Resume-Matching-Tool-for-Job-Descriptions-recruitRyte.jpg"
CHATBOT_IMAGE_URL = "https://neo4j.com/wp-content/uploads/2023/04/knowledge-graph-based-chatbot-scaled.jpg"





st.set_page_config(page_title="AI Career Assistant", layout="wide", page_icon="💼")
st.markdown("""
<style>
    /* General Streamlit Overrides */
    body { background-color: #0e1117; color: #ffffff; }
    .stButton>button { background-color: #262730; color: white; border-radius: 10px; padding: 0.6em 1em; border: none; }
    .stButton>button:hover { background-color: #333642; }
    .stCard { 
        border: 1px solid #333642; 
        border-radius: 12px; 
        padding: 10px; 
        background-color: #1e1e2d; 
        margin-top: 15px;
    }
   
    /* Custom style for st.expander in Resume Matcher */
    .stExpanderHeader {
        background-color: #1e1e2d; /* Darker background for headers */
        border-radius: 8px;
        padding: 10px;
        margin-bottom: 5px;
        font-weight: bold;
        border: 1px solid #333642;
    }
    
    /* ****************************************************** */
    /* *** MODIFICATION: REMOVE TOP GAP AND CENTER HEADERS *** */
    /* ****************************************************** */
    /* Target the main content container and reduce its top padding/margin */
    .block-container {
        padding-top: 2rem; /* Reduced from Streamlit default (e.g., 5rem) */
        padding-bottom: 2rem;
        padding-left: 2rem;
        padding-right: 2rem;
    }

    /* Target the main Streamlit h1/h2 tags for centering */
    .main-centered-title {
        text-align: center;
        margin-top: 0 !important; /* Eliminate top margin for the centered titles */
        margin-bottom: 1.5rem;
        font-size: 2.5em;
        color: #f0f0f5;
    }
    .st-emotion-cache-18ni7ap { /* Target for the main component wrapper */
        padding-top: 1rem;
    }
</style>
""", unsafe_allow_html=True)


BASE_PATH = r"/content/drive/MyDrive/gohak/Rag_model_and_chatbot"
ASSETS_PATH = os.path.join(BASE_PATH, "assets")
SINGLE_BANNER_PATH = os.path.join(ASSETS_PATH, "banner7.png")

def get_image_as_bytes(path):
    """Reads a local image file into bytes."""
    try:

        if os.path.exists(path):
            with open(path, "rb") as image_file:
                return image_file.read()
        else:
            return None
    except Exception:
        return None


@st.cache_resource(show_spinner=False)
def load_resume_model():
    from sentence_transformers import SentenceTransformer
    try:
        model = SentenceTransformer("all-MiniLM-L6-v2", device="cpu")
        return model
    except Exception as e:
        st.error(f"❌ Resume matching model failed to load: {e}")
        return None

@st.cache_resource(show_spinner=False)
def load_flan_llm():
    model_name = "google/flan-t5-small"
    try:
        tokenizer = AutoTokenizer.from_pretrained(model_name)
        model = AutoModelForSeq2SeqLM.from_pretrained(model_name)
      
        pipe = pipeline("text2text-generation", model=model, tokenizer=tokenizer, device=-1, max_new_tokens=200)
        return HuggingFacePipeline(pipeline=pipe)
    except Exception:
        return None

@st.cache_resource(show_spinner=False)
def load_embeddings_and_faiss():
    try:
        emb_model = HuggingFaceEmbeddings(model_name="sentence-transformers/paraphrase-MiniLM-L3-v2")
     
        faiss_path = os.path.join(BASE_PATH, "faiss_index")
        if not os.path.exists(faiss_path):
            return emb_model, None
        
        vect = FAISS.load_local(faiss_path, embeddings=emb_model, allow_dangerous_deserialization=True)
        return emb_model, vect
    except Exception:
     
        emb_model = HuggingFaceEmbeddings(model_name="sentence-transformers/paraphrase-MiniLM-L3-v2")
        return emb_model, None


@st.cache_resource(show_spinner=False)
def build_rag_chain():
    from langchain_core.prompts import PromptTemplate
    from langchain_core.runnables import RunnablePassthrough
    from langchain_core.output_parsers import StrOutputParser
    from langchain_groq import ChatGroq

    _, vectorstore = load_embeddings_and_faiss()

    if vectorstore is None:
        st.error("❌ FAISS vector store not found.")
        return None

    # LLM
    llm = ChatGroq(
        groq_api_key="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",
        model_name="llama-3.1-8b-instant",
        temperature=0.2
    )

    prompt = PromptTemplate.from_template("""
You are an expert IT career assistant.
Answer clearly and professionally using ONLY the context.
If context does not contain the answer, say "I don't have that exact information."

Context:
{context}

Question:
{question}

Final Answer:
""")

    retriever = vectorstore.as_retriever(search_kwargs={"k": 4})

    # Build the LCEL RAG chain (Latest method, NO RetrievalQA)
    rag_chain = (
        {"context": retriever, "question": RunnablePassthrough()}
        | prompt
        | llm
        | StrOutputParser()
    )

    return rag_chain



@st.cache_data
def load_data():
    try:
        
        dataset = load_dataset("DevilsLord/It_job_roles_skills_certifications")
        df = dataset['train'].to_pandas()
       
        df['combined'] = df['Job Description'].fillna('') + " " + df['Skills'].fillna('') + " " + df['Certifications'].fillna('')
        return df
    except Exception:
    
        return pd.DataFrame({'Job Title': [], 'Job Description': [], 'Skills': [], 'Certifications': [], 'combined': []})


with st.spinner("Initializing AI models and data..."):
    df = load_data()
    resume_model = load_resume_model()

    #

    rag_chain = build_rag_chain()
   


def extract_text(file):
    ext = file.name.split('.')[-1].lower()
    if ext == 'txt': return file.read().decode('utf-8', errors='ignore')
    elif ext == 'pdf':
        try:
            
            doc = fitz.open(stream=file.read(), filetype="pdf")
            return "\n".join([page.get_text("text") for page in doc])
        except Exception:
            return ""
    elif ext == 'docx': 
        try:
            return docx2txt.process(file)
        except Exception:
            return ""
    elif ext == 'pptx':
        text = ""
        try:
            presentation = Presentation(file)
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"): text += shape.text + "\n"
        except Exception:
            return ""
        return text
    return ""


st.sidebar.title("🌙 Navigation")
page = st.sidebar.radio("Go to", ["🏠 Home", "📄 Resume Matcher", "💬 Chatbot"], label_visibility="collapsed")


if page == "🏠 Home":
   
    st.markdown('<h1 class="main-centered-title">🧠 AI-Powered Career Assistant</h1>', unsafe_allow_html=True)
    st.markdown("<p style='text-align: center; margin-top: -15px;'>Welcome to your <b>AI-driven Resume–Job Matcher & Chatbot</b>!</p>", unsafe_allow_html=True)

    col1, col2 = st.columns([3, 1])

    with col1:
    
        image_bytes = get_image_as_bytes(SINGLE_BANNER_PATH)
        st.markdown('<div class="banner-image-wrapper">', unsafe_allow_html=True)
        if image_bytes:
            st.image(image_bytes, use_container_width=True, caption="AI Career Assistant Banner")
        else:
            st.warning(f"Image not found at: {SINGLE_BANNER_PATH}. Please ensure banner1.png exists.")
            st.markdown(
                '<div style="height: 100%; display: flex; align-items: center; justify-content: center; color: #888;">'
                'Image Placeholder: banner1.png not loaded.'
                '</div>', unsafe_allow_html=True
            )
        st.markdown('</div>', unsafe_allow_html=True)
        
     
        st.markdown("### 🎯 Our Mission")
        st.markdown("""
        In today's fast-paced tech landscape, finding the perfect job or the right talent can be challenging. 
        Our AI Career Assistant uses **advanced Natural Language Processing (NLP)** and **Machine Learning (ML)** to bridge the gap between job seekers and IT opportunities.
        
        We aim to provide precise career guidance by matching your unique skills from your resume with thousands 
        of verified IT job descriptions, and by offering instant, knowledgeable answers through our RAG-based chatbot.
        """)
        st.info("Upload your resume now or ask the chatbot about in-demand skills and certifications to get started!")
        
        st.markdown("---")
        
        
        st.markdown("### 🚀 Core Features")
        st.markdown("""
<div style="padding-left: 15px;">
    <ul style="list-style-type: '👉';">
        <li>
            <b>Semantic Resume-Job Matching:</b> Utilizes advanced Sentence-BERT models to analyze the *meaning* of your skills, not just keywords, ensuring highly accurate matches with over 200 IT job roles.
        </li>
        <li>
            <b>Retrieval-Augmented Generation (RAG) Chatbot:</b> A specialized Q&A system built on FAISS and Flan-T5, providing context-aware, up-to-date answers on IT skills, certifications, and career paths, drawing directly from a verified knowledge base.
        </li>
        <li>
            <b>Multi-Format Resume Parsing:</b> Seamlessly extracts text from various document types including PDF, DOCX, PPTX, and TXT, making the analysis process quick and accessible.
        </li>
        <li>
            <b>Intuitive Dark UI/UX:</b> Features a sleek, responsive dark-themed interface with tabbed navigation (Resume Matcher and Chatbot) for a premium and efficient user experience.
        </li>
    </ul>
</div>
""", unsafe_allow_html=True)
        st.markdown("---")
        st.info("Navigate to the **Resume Matcher** or **Chatbot** using the menu on the left to start accelerating your career!")

   
    with col2:
        
        st.markdown('<div class="stCard">', unsafe_allow_html=True)
        st.subheader("💡 Quick Info")
        
        if df is not None:
            st.markdown(f"**Total Job Roles:** `{len(df)}`")
        else:
            st.markdown(f"**Total Job Roles:** `Data Error`")
            
        st.markdown("**Core Models:**")
        st.markdown("- **Matcher:** Sentence-BERT")
        st.markdown("- **Chatbot:** FLAN-T5 Small (RAG)")
        st.markdown("**Knowledge Base:**")
        st.markdown("- IT Job Descriptions")
        st.markdown("- Required Skills & Certs")
        st.markdown("</div>", unsafe_allow_html=True)
        
        st.markdown("---")
        st.subheader("⚙️ How it Works")
        st.write("The system uses deep learning models (Sentence Transformers) for semantic matching and a Retrieval-Augmented Generation (RAG) system based on FAISS and Flan-T5 for Q&A.")



elif page == "📄 Resume Matcher":
    
    st.markdown('<h2 class="main-centered-title">📄 AI Resume Matcher</h2>', unsafe_allow_html=True)
    st.markdown("<p style='text-align: center; margin-top: -15px;'>Find your perfect IT career path by matching your skills against our extensive database of job roles.</p>", unsafe_allow_html=True)
    
   
    st.markdown("---")


    upload_col, info_col = st.columns([1.5, 2])
    
    with upload_col:
        st.subheader("1. Upload Your Resume")
        uploaded_file = st.file_uploader("PDF, DOCX, PPTX, TXT", type=["txt", "pdf", "docx", "pptx"])
        
    with info_col:
        st.subheader("2. How Matching Works")
        st.info("""
        We use an AI model (**Sentence-BERT**) to analyze the key skills and experience in your resume 
        and compare them semantically with the requirements of hundreds of IT job roles. 
        
        The result is a **Match Score (0.0 to 1.0)**, where a score closer to 1.0 indicates a higher alignment with the job's requirements.
        """)
        
    if uploaded_file:
        resume_text = extract_text(uploaded_file).replace('\n',' ').strip()
        
        if resume_model is None or df is None:
            st.error("❌ Model or data is not loaded. Cannot perform matching.")
        elif not resume_text:
            st.error("❌ No readable text found in the uploaded file. Please ensure your file is not corrupted or password-protected.")
        elif resume_text:
            st.divider()
            st.subheader("✨ Top Job Matches")
            
            with st.spinner("⏳ Analyzing resume and finding best job matches..."):
                
                if resume_text:
                    resume_emb = resume_model.encode(resume_text, convert_to_tensor=True)
                    job_embs = resume_model.encode(df['combined'].tolist(), convert_to_tensor=True)
                    sims = util.cos_sim(resume_emb, job_embs)[0].cpu().numpy()
                    df['Similarity'] = sims
                    top_matches = df.sort_values('Similarity', ascending=False).head(10).reset_index(drop=True)
                else:
                    top_matches = pd.DataFrame()

            if not top_matches.empty:
                for index, row in top_matches.iterrows():
                    match_score = row['Similarity']
                    
                    if match_score >= 0.8:
                        color = "green"
                        icon = "🔥"
                    elif match_score >= 0.6:
                        color = "orange"
                        icon = "⭐"
                    else:
                        color = "yellow"
                        icon = "💡"
                        
                    
                    expander_label = f"{icon} **{row['Job Title']}** — Match Score: **{match_score:.2f}**"
                    
                    with st.expander(expander_label, expanded=False):
                        st.markdown("#### Job Details")
                        st.markdown(f"**📝 Description:** {row['Job Description']}")
                        st.markdown("---")
                        st.markdown(f"**🧩 Key Skills:** {row['Skills']}")
                        st.markdown("---")
                        st.markdown(f"**📜 Certifications:** {row['Certifications']}")
            else:
                st.warning("No job roles found or an error occurred during matching.")



elif page == "💬 Chatbot":
   
    st.markdown('<h2 class="main-centered-title">💬 Custom IT Career Chatbot</h2>', unsafe_allow_html=True)
    st.markdown("<p style='text-align: center; margin-top: -15px;'>Ask me anything about IT job roles, required skills, or popular certifications.</p>", unsafe_allow_html=True)

   
    st.markdown("---")
    
    if rag_chain is None:
        st.error("❌ Chatbot model (RAG chain) is not loaded. Please check model paths and try reloading.")
        st.stop()
        
   
    if "messages" not in st.session_state:
        st.session_state["messages"] = [{"role": "assistant", "content": "Hello! I'm your AI career assistant. What IT career question can I answer for you today?"}]

   
    for msg in st.session_state.messages:
        with st.chat_message(msg["role"]):
            st.write(msg["content"])

   
    if prompt := st.chat_input("Ask a question about IT roles, skills, or certifications..."):
      
        st.session_state.messages.append({"role": "user", "content": prompt})
        with st.chat_message("user"):
            st.write(prompt)

        
        with st.chat_message("assistant"):
            with st.spinner("🤔 Thinking and retrieving knowledge..."):
                assistant_response = rag_chain.invoke(prompt)


           
            
          
            st.markdown("### ✅ Answer")
            st.write(assistant_response)
            
            
            

           
            st.session_state.messages.append({"role": "assistant", "content": assistant_response})
           


def main():
  
    pass 

if __name__ == "__main__":
    main()